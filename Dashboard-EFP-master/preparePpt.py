from xlrd import open_workbook
from collections import OrderedDict
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_TICK_LABEL_POSITION
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_TICK_MARK
from pptx.spec import autoshape_types
from pptx.enum.shapes import MSO_SHAPE



from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
import csv
import sys
import os
from pptx.api import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
import datetime
from time import strftime, localtime
from pptx.dml.color import ColorFormat, RGBColor
import xlrd
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_TICK_LABEL_POSITION
import re
import collections
import csv
from calendar import calendar
img_path="./../../../../Hotel fILES/Strome Business Data/TopBanner.png"

header_img="./../../../../Hotel fILES/Strome Business Data/header.png"

prs = Presentation()


slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.shapes.add_picture(header_img, Inches(0), Inches(0))

left = top =  Inches(2)
width = height  =  Inches(5)
txBox = slide.shapes.add_textbox(Inches(0), top ,Inches(10),Inches(10))
tf = txBox.text_frame

p = tf.add_paragraph()
p.text = "CEAP Data Analysis"
p.font.bold = True
p.alignment =  PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = ""

p = tf.add_paragraph()
p.text = strftime("%b %d %Y", localtime())
p.font.bold = True
p.alignment =  PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = ""
p = tf.add_paragraph()
p.text = ""

p = tf.add_paragraph()
p.text = "Professor Vinod Agarwal"
p.alignment =  PP_ALIGN.CENTER


p = tf.add_paragraph()
p.text = ""

p = tf.add_paragraph()
p.text = "www.odu.edu/forecasting"
p.font.underline = True
p.alignment =  PP_ALIGN.CENTER






######################################         TNF SSA         ########################################


# title_slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(title_slide_layout)
# title = slide.shapes.title
# subtitle = slide.placeholders[1]
#
# title.text = "Seasonally Adjusted Nonfarm Employment by Month"
# subtitle.text = " "

regionNames={"Virginia":"Virginia",
"Virginia Beach-Norfolk-Newport News, VA-NC MSA":"VirginiaBeach",
"Washington-Arlington-Alexandria, DC-VA-MD-WV MSA, VA part":"Washington",
"Blacksburg-Christiansburg-Radford, VA MSA":"Blacksburg",
"Charlottesville, VA MSA":"Charlottesville",
"Winchester, VA-WV MSA":"Winchester",
"Harrisonburg, VA MSA":"Harrisonburg",
"Lynchburg, VA MSA":"Lynchburg",
"Richmond, VA MSA":"Richmond",
"Roanoke, VA MSA":"Roanoke",
"Harrisonburg, VA MSA":"Harrisonburg"};


f = open('./Csv/EMPLOYEMENT_TNFALL_emp_Monthly.csv')
csv_f = csv.reader(f)
a=[]

for row in csv_f:
  a.append(row)
a=zip(*a)
#print a
b=[]
max1 = 0
min1 = 0
min1 = a[0][1].split('/')[-1]
max1 = a[0][-1].split('/')[-1]
# print a[7]
for k in range(1,len(a)):
    if 1:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(header_img, Inches(0), Inches(0))
        chart_data = ChartData()
        b=[]
        for m in range(0,len(a[0])):
            if m!=0:
                b.append(a[0][m])
        chart_data.categories = b
        row=[]
        for l in range(1,len(a[k])):
            if l==len(a[k]) or l==len(a[k])-1 or l==len(a[k])-2:
                if float(a[k][l])!=0:
                    if a[k][0]!='US':
                        row.append(float(a[k][l])/1000)
            else:
                if a[k][0]!='US':
                    row.append(float(a[k][l])/1000)
            # if float(a[k][l])/1000 > max :
            #     max = float(a[k][l])/1000
            # if float(a[k][l])/1000 < min :
            #     min = float(a[k][l])/1000
        txBox = slide.shapes.add_textbox(Inches(0), Inches(0.75) ,Inches(10),Inches(0.5))
        tf = txBox.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        run = p1.add_run()
        oo=a[k][0]
        run.text ="Seasonally Adjusted Nonfarm Employment by Month for " + regionNames[oo] + " (in Thousands)"
        domainStart = {  "Virginia" :    2500,  
              "VirginiaBeach" :    700,   
              "Blacksburg" :    65,   
              "Charlottesville" :    90,    
              "Harrisonburg" :    55,   
              "Lynchburg" :    90,   
              "Richmond" :    550,  
              "Roanoke" :    120, 
              "Staunton" :    40, 
              "Winchester" :    50, 
              "Washington" :    1000 }


        font=run.font
        font.size = Pt(16)
        font.bold = True
        font.color.theme_color = MSO_THEME_COLOR.ACCENT_2

        p = tf.add_paragraph()
        run = p.add_run()
        run.text = min1 + "-" + max1

        shape = slide.shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_1, Inches(0.1), Inches(1.05) ,Inches(6),Inches(0.02))
        line = shape.line
        line.fill.background()

        font = run.font
        font.name = 'Calibri'
        font.size = Pt(16)
        font.bold = True
        font.italic = None  # cause value to be inherited from theme
        font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
        #print a[k][0],row
        #print a[k][0]
        chart_data.add_series(a[k][0],row)
        chart_data.add_series(a[k][0],'0')
        x, y, cx, cy = Inches(0), Inches(1.5), Inches(10), Inches(5.75)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
        ).chart
        chart.has_legend = False
        #chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        #chart.legend.include_in_layout = False
        chart.series[0].smooth = True
        #chart.has_legend = True
        value_axis = chart.value_axis
        value_axis.has_title = True
        category_axis = chart.category_axis
        value_axis.tick_labels.font.bold = True
        value_axis.tick_labels.font.size = Pt(12)
        category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
        category_axis.tick_labels.font.bold = True
        category_axis.tick_labels.font.size = Pt(12)
        category_axis.tick_labels.offset =10
        tick_labels = value_axis.tick_labels
        tick_labels.number_format = '0.00'
        #chart.value_axis.maximum_scale = 4000.00
        if regionNames[oo] == 'Roanoke':
            chart.value_axis.minimum_scale = 140
        elif regionNames[oo] == 'Washington':
            chart.value_axis.minimum_scale = 2500
        else:    
            chart.value_axis.minimum_scale = domainStart[regionNames[oo]]   


        # chart.has_legend = False
        # chart.series[0].smooth = True
        # value_axis = chart.value_axis
        # tick_labels = value_axis.tick_labels
        # tick_labels.number_format = '0.00'
        # tick_labels.font.bold = True
        # tick_labels.font.size = Pt(16)
        # category_axis=chart.category_axis
        # tick_labels = category_axis.tick_labels
        # tick_labels.font.bold = True
        # tick_labels.font.size = Pt(16)
        #
        # fill = chart.series[0].format.fill
        # fill.solid()
        # fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_5








######################################         TNF NSA         ########################################




# title_slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(title_slide_layout)
# title = slide.shapes.title
# subtitle = slide.placeholders[1]
#
# title.text = "Not Seasonally Adjusted Nonfarm Employment by Month"
# subtitle.text = " "

regionNames={"Virginia":"Virginia",
"Virginia Beach-Norfolk-Newport News, VA-NC MSA":"VirginiaBeach",
"Washington-Arlington-Alexandria, DC-VA-MD-WV MSA, VA part":"Washington",
"Blacksburg-Christiansburg-Radford, VA MSA":"Blacksburg",
"Charlottesville, VA MSA":"Charlottesville",
"Winchester, VA-WV MSA":"Winchester",
"Harrisonburg, VA MSA":"Harrisonburg",
"Lynchburg, VA MSA":"Lynchburg",
"Richmond, VA MSA":"Richmond",
"Roanoke, VA MSA":"Roanoke",
"Staunton-Waynesboro, VA":"Staunton"};


f1 = open('./Csv/EMPLOYEMENT_TNFALLNSA_emp_Monthly2.csv')
csv_f1 = csv.reader(f1)
a1=[]

for row in csv_f1:
  a1.append(row)
a1=zip(*a1)
#print a1
b1=[]

max1 = 0
min1 = 0
min1 = a1[0][1].split('/')[-1]
max1 = a1[0][-1].split('/')[-1]

for k in range(1,len(a1)):

    if a1[k][0]!='US' and a1[k][0]!='NOVA':
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(header_img, Inches(0), Inches(0))

        #slide = prs.slides.add_slide(prs.slide_layouts[5])
        #title = slide.shapes.title
        chart_data = ChartData()
        b=[]
        for m in range(0,len(a1[0])):
            if m!=0:
                if a1[0][m]!='US':
                    b1.append(a1[0][m])
        chart_data.categories = b1
        row=[]
        for l in range(1,len(a1[k])):
            if l==len(a1[k]) or l==len(a1[k])-1 or l==len(a1[k])-2:
                if float(a1[k][l])!=0:
                    if a1[k][0]!='US':
                        row.append(float(a1[k][l])/1000)
            else:
                if a1[k][0]!='US':
                    row.append(float(a1[k][l])/1000)


        txBox = slide.shapes.add_textbox(Inches(0), Inches(0.75) ,Inches(10),Inches(0.5))
        tf = txBox.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        run = p1.add_run()
        #print a1[k]
        oo=a1[k][0]
        run.text ="Not Seasonally Adjusted Nonfarm Employment by Month for " + regionNames[oo]+ " (in Thousands)"
        domainStart = {  "Virginia" :    2500,  
              "VirginiaBeach" :    700,   
              "Blacksburg" :    65,   
              "Charlottesville" :    90,    
              "Harrisonburg" :    55,   
              "Lynchburg" :    90,   
              "Richmond" :    550,  
              "Roanoke" :    120, 
              "Staunton" :    40, 
              "Winchester" :    50, 
              "Washington" :    1000 }
        font=run.font
        font.size = Pt(16)
        font.bold = True
        font.color.theme_color = MSO_THEME_COLOR.ACCENT_2

        p = tf.add_paragraph()
        run = p.add_run()
        run.text = min1 + "-" + max1

        shape = slide.shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_1, Inches(0.1), Inches(1.05) ,Inches(6),Inches(0.02))
        line = shape.line
        line.fill.background()

        font = run.font
        font.name = 'Calibri'
        font.size = Pt(16)
        font.bold = True
        font.italic = None  # cause value to be inherited from theme
        font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
        if 1:
            #if regionNames[oo] == 'Washington':
                #print row
            chart_data.add_series(a1[k][0],row)
            chart_data.add_series(a1[k][0],'0')
            x, y, cx, cy = Inches(0), Inches(1.5), Inches(10), Inches(5.75)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
            ).chart

            chart.has_legend = False
            chart.series[0].smooth = True
            value_axis = chart.value_axis
            tick_labels = value_axis.tick_labels
            tick_labels.number_format = '0.00'
            tick_labels.font.bold = True
            tick_labels.font.size = Pt(12)
            category_axis=chart.category_axis
            tick_labels = category_axis.tick_labels
            tick_labels.font.bold = True
            tick_labels.font.size = Pt(12)
            if regionNames[oo] == 'Roanoke':
                chart.value_axis.minimum_scale = 140
            elif regionNames[oo] == 'Washington':
                chart.value_axis.minimum_scale = 2500
            else:
                chart.value_axis.minimum_scale = domainStart[regionNames[oo]]   
            #chart.value_axis.minimum_scale = domainStart[regionNames[oo]]


##################################################        RGDP quartelrly VA                         ######################################

# title_slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(title_slide_layout)
# title = slide.shapes.title
# subtitle = slide.placeholders[1]
#
# title.text = "RGDP Quarterly VA"
# subtitle.text = " "

f2 = open("./Csv/VA_gdp.csv",'r')

csv_f2 = csv.reader(f2)
a1=[]
for row in csv_f2:
  a1.append(row)
a1=zip(*a1)

b1=[]

max1 = 0
min1 = 0
min1 = a1[0][1].split('/')[-1]
max1 = a1[0][-1].split('/')[-1]

for m in range(1,len(a1[0])):
    b1.append(a1[0][m])
for l in range(len(a1)-2,len(a1)-1):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(header_img, Inches(0), Inches(0))


    row=[]
    oo = a1[l][0]
    for k in range(1,len(a1[l])):
        # if k==len(a1[l])-1:
        #     if a1[l][k]=='0':
        #         print b1
        #         print "value missing"
        #         b1 = b1[:-1]
        #         print b1
        #     else:
        #         row.append(float(a1[l][k])/1000)
        # else:
        row.append(float(a1[l][k])/1000)

    txBox = slide.shapes.add_textbox(Inches(0), Inches(0.75) ,Inches(10),Inches(0.5))
    tf = txBox.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    run = p1.add_run()
    #$oo=a[k][0].split('-')[0]
    run.text ="Real Gross Domestic Product by Quarter for Virginia in Billions of Dollars"
    font=run.font
    font.size = Pt(16)
    font.bold = True
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_2

    p = tf.add_paragraph()
    run = p.add_run()
    run.text = min1 + "-" + max1

    shape = slide.shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_1, Inches(0.1), Inches(1.05) ,Inches(4),Inches(0.02))
    line = shape.line
    line.fill.background()

    font = run.font
    font.name = 'Calibri'
    font.size = Pt(16)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

    chart_data = ChartData()
    chart_data.categories = b1

    chart_data.add_series(a1[l][0],row)
    #title = slide.shapes.title
    #title.text = a1[l][0]
    #chart_data.add_series('Series 1', (19.2, 21.4, 16.7))

    # add chart to slide --------------------
    x, y, cx, cy = Inches(0), Inches(1.5), Inches(10), Inches(5.75)
    chart=slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    value_series=chart.series[0]
    value_series.fill.solid()
    value_series.fill.fore_color.theme_color=MSO_THEME_COLOR.ACCENT_5
    chart.has_legend = False
    chart.series[0].smooth = True
    value_axis = chart.value_axis
    tick_labels = value_axis.tick_labels
    tick_labels.number_format = '"$"0.00'
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(12)
    category_axis=chart.category_axis
    tick_labels = category_axis.tick_labels
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(12)
    chart.value_axis.minimum_scale = 0

    # value_series=chart.series[0]
    # value_series.fill.solid()
    # value_series.fill.fore_color.theme_color=MSO_THEME_COLOR.ACCENT_5
    # chart.has_legend = False
    # chart.series[0].smooth = True
    # value_axis = chart.value_axis
    # tick_labels = value_axis.tick_labels
    # tick_labels.number_format = '#,###'
    # tick_labels.font.bold = True
    # tick_labels.font.size = Pt(11)
    # category_axis=chart.category_axis
    # tick_labels = category_axis.tick_labels
    # tick_labels.font.bold = True
    # tick_labels.font.size = Pt(11)


##################################################        Analysed VA gdp                         ######################################

# title_slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(title_slide_layout)
# title = slide.shapes.title
# subtitle = slide.placeholders[1]
#
# title.text = "RGDP Quarterly VA"
# subtitle.text = " "

f2 = open("./Csv/VA_gdp.csv",'r')

csv_f2 = csv.reader(f2)
a1=[]
for row in csv_f2:
  a1.append(row)
a1=zip(*a1)
#print a1
b1=[]

max1 = 0
min1 = 0
min1 = a1[0][2].split('/')[-1]
max1 = a1[0][-1].split('/')[-1]

for m in range(2,len(a1[0])):
    b1.append(a1[0][m])
for l in range(len(a1)-1,len(a1)):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(header_img, Inches(0), Inches(0))


    row=[]
    oo = a1[l][0]
    # print a1
    # break
    for k in range(2,len(a1[l])):
        if float(a1[l][k])!=-400:
            row.append(float(a1[l][k]))

    txBox = slide.shapes.add_textbox(Inches(0), Inches(0.75) ,Inches(10),Inches(0.5))
    tf = txBox.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    run = p1.add_run()
    #$oo=a[k][0].split('-')[0]
    run.text ="Analyzed Gross Domestic Product by Quarter for Virginia"
    font=run.font
    font.size = Pt(16)
    font.bold = True
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_2

    p = tf.add_paragraph()
    run = p.add_run()
    run.text = min1 + "-" + max1

    shape = slide.shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_1, Inches(0.1), Inches(1.05) ,Inches(4),Inches(0.02))
    line = shape.line
    line.fill.background()

    font = run.font
    font.name = 'Calibri'
    font.size = Pt(16)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

    chart_data = ChartData()
    chart_data.categories = b1

    chart_data.add_series(a1[l][0],row)
    #title = slide.shapes.title
    #title.text = a1[l][0]
    #chart_data.add_series('Series 1', (19.2, 21.4, 16.7))

    # add chart to slide --------------------
    x, y, cx, cy = Inches(0), Inches(1.5), Inches(10), Inches(5.75)
    chart=slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    value_series=chart.series[0]
    value_series.fill.solid()
    value_series.fill.fore_color.theme_color=MSO_THEME_COLOR.ACCENT_3
    chart.has_legend = False
    chart.series[0].smooth = True
    chart.series[0].invert_if_negative = False
    category_axis = chart.category_axis
    category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
    value_axis = chart.value_axis
    tick_labels = value_axis.tick_labels
    tick_labels.number_format = '0.0"%"'
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(12)
    category_axis=chart.category_axis
    tick_labels = category_axis.tick_labels
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(12)

    # value_series=chart.series[0]
    # value_series.fill.solid()
    # value_series.fill.fore_color.theme_color=MSO_THEME_COLOR.ACCENT_5
    # chart.has_legend = False
    # chart.series[0].smooth = True
    # value_axis = chart.value_axis
    # tick_labels = value_axis.tick_labels
    # tick_labels.number_format = '#,###'
    # tick_labels.font.bold = True
    # tick_labels.font.size = Pt(11)
    # category_axis=chart.category_axis
    # tick_labels = category_axis.tick_labels
    # tick_labels.font.bold = True
    # tick_labels.font.size = Pt(11)





##################################################        RGDP                         ######################################

# title_slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(title_slide_layout)
# title = slide.shapes.title
# subtitle = slide.placeholders[1]
#
# title.text = "RGDP"
# subtitle.text = " "

f2 = open("./Csv/RGDP.csv",'r')

csv_f2 = csv.reader(f2)
a1=[]
for row in csv_f2:
  a1.append(row)
a1=zip(*a1)

b1=[]

max1 = 0
min1 = 0
min1 = a1[0][1]
max1 = a1[0][-1]

for m in range(1,len(a1[0])):
    b1.append(a1[0][m])
print b1
flag=0
for l in range(1,len(a1)):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(header_img, Inches(0), Inches(0))


    row=[]
    oo = a1[l][0]
    print a1[l]
    if oo != "Virginia":
        if flag==0:
            b1.remove('2016')
            flag=1
    else:
        b1.append('2016')

    for k in range(1,len(a1[l])):

        if k==len(a1[l])-1:
            if a1[l][k]=='0':
                print "value missing"
            else:
                row.append(float(a1[l][k])/1000)
        else:
            row.append(float(a1[l][k])/1000)

    txBox = slide.shapes.add_textbox(Inches(0), Inches(0.75) ,Inches(10),Inches(0.5))
    tf = txBox.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    run = p1.add_run()
    #$oo=a[k][0].split('-')[0]
    run.text ="Real Gross Domestic Product by Year for " + oo + " in Billions of Dollars"
    font=run.font
    font.size = Pt(16)
    font.bold = True
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_2

    p = tf.add_paragraph()
    run = p.add_run()
    run.text = min1 + "-" + max1

    shape = slide.shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_1, Inches(0.1), Inches(1.05) ,Inches(4),Inches(0.02))
    line = shape.line
    line.fill.background()

    font = run.font
    font.name = 'Calibri'
    font.size = Pt(16)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

    chart_data = ChartData()
    chart_data.categories = b1

    chart_data.add_series(a1[l][0],row)
    #title = slide.shapes.title
    #title.text = a1[l][0]
    #chart_data.add_series('Series 1', (19.2, 21.4, 16.7))

    # add chart to slide --------------------
    x, y, cx, cy = Inches(0), Inches(1.5), Inches(10), Inches(5.75)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    value_series=chart.series[0]
    value_series.fill.solid()
    value_series.fill.fore_color.theme_color=MSO_THEME_COLOR.ACCENT_5

    chart.has_legend = False
    chart.series[0].smooth = True
    value_axis = chart.value_axis
    tick_labels = value_axis.tick_labels
    tick_labels.number_format = '"$"#0'
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(12)
    category_axis=chart.category_axis
    tick_labels = category_axis.tick_labels
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(12)
    chart.value_axis.minimum_scale = 0


##################################################        Growth Ratee RGDP                         ######################################

# title_slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(title_slide_layout)
# title = slide.shapes.title
# subtitle = slide.placeholders[1]
#
# title.text = "Grwoth Rate RGDP"
# subtitle.text = " "

f15 = open("./Csv/RGDP.csv",'r')

csv_f15 = csv.reader(f15)
a1=[]
for row in csv_f15:
  a1.append(row)
a1=zip(*a1)

b1=[]

max1 = 0
min1 = 0
min1 = a1[0][2]
max1 = a1[0][-1]
flag=0
for m in range(2,len(a1[0])):
    b1.append(a1[0][m])
for l in range(1,len(a1)):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(header_img, Inches(0), Inches(0))



    row=[]
    oo = a1[l][0]

    if oo!="Virginia":
        if flag==0:
            b1.remove('2016')
            flag=1
    else:
        b1.append('2016')

        
    for k in range(2,len(a1[l])):
        if(float(a1[l][k])!=0):
            diff = (float(a1[l][k])-float(a1[l][k-1]))/float(a1[l][k-1]) * 100
            row.append(diff)

    txBox = slide.shapes.add_textbox(Inches(0), Inches(0.75) ,Inches(10),Inches(0.5))
    tf = txBox.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    run = p1.add_run()
    #$oo=a[k][0].split('-')[0]tnf
    run.text ="Growth Rate of Real Gross Domestic Product by Year for " + oo
    font=run.font
    font.size = Pt(16)
    font.bold = True
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_2

    p = tf.add_paragraph()
    run = p.add_run()
    run.text = min1 + "-" + max1

    shape = slide.shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_1, Inches(0.1), Inches(1.05) ,Inches(4),Inches(0.02))
    line = shape.line
    line.fill.background()

    font = run.font
    font.name = 'Calibri'
    font.size = Pt(16)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

    chart_data = ChartData()
    chart_data.categories = b1

    chart_data.add_series(a1[l][0],row)
    #title = slide.shapes.title
    #title.text = a1[l][0]
    #chart_data.add_series('Series 1', (19.2, 21.4, 16.7))

    # add chart to slide --------------------



    x, y, cx, cy = Inches(0), Inches(1.5), Inches(10), Inches(5.75)
    chart=slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    value_series=chart.series[0]
    value_series.fill.solid()
    value_series.fill.fore_color.theme_color=MSO_THEME_COLOR.ACCENT_3
    chart.has_legend = False
    chart.series[0].smooth = True
    chart.series[0].invert_if_negative = False
    category_axis = chart.category_axis
    category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
    value_axis = chart.value_axis
    tick_labels = value_axis.tick_labels
    tick_labels.number_format = '0.0"%"'
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(12)
    category_axis=chart.category_axis
    tick_labels = category_axis.tick_labels
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(12)







##################################################        REVPAR                         ######################################

# title_slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(title_slide_layout)
# title = slide.shapes.title
# subtitle = slide.placeholders[1]
# title.text = "Hotel Revenue per Available Room by Month for"
# subtitle.text = " "


f3 = open("./Csv/Dashboard_REVPAR_Month1.csv",'r')

dic1={"WashingtonDC-MD-VAMarket":"Washington DC",
     "Richmond/PetersburgMarket":"Richmond/Petersburg",
     "Staunton/HarrisonburgMarket":"Staunton/Harrisonburg",
     "VirginiaPortionofWashingtonDC":"Virginia",
     "Blacksburg/WythevilleMarket":"Blacksburg",
     "theCommonwealthofVirginia":"Commonwealth of Virginia",
     "CharlottesvilleMarket":"Charlottesville",
     "LynchburgMarket":"Lynchburg",
     "RoanokeMarket":"Roanoke",
     "HamptonRoadsMarket":"Hampton Roads"}

dic={"Jan":"01","Feb":"02","Mar":"03","Apr":"04","May":"05","Jun":"06","Jul":"07","Aug":"08","Sep":"09","Oct":"10","Nov":"11","Dec":"12"}
csv_f3 = csv.reader(f3)
a1=[]
for row in csv_f3:
  a1.append(row)
a1=zip(*a1)
b=[]

max1 = 0
min1 = 0
min1 = a1[0][1]
max1 = a1[0][-1]


for row_num in range(1,len(a1[0])):
    b.append(dic[a1[1][row_num]]+"-01-"+a1[0][row_num])
a2=[]
a2.append(b)
#print a2
for row in range(2,len(a1)):
    a2.append(a1[row])
#print a2
for k in range(1,len(a2)):

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(header_img, Inches(0), Inches(0))


    b1=[]
    for m in range(0,len(a2[0])):
        b1.append(a2[0][m])

    row=[]
    #print a2[k]
    for l in range(1,len(a2[k])):
        #print b[k]
        row.append(float(a2[k][l]))

    txBox = slide.shapes.add_textbox(Inches(0), Inches(0.75) ,Inches(10),Inches(0.5))
    tf = txBox.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    run = p1.add_run()
    #$oo=a[k][0].split('-')[0]
    run.text ="Hotel Revenue per Available Room by Month for " + dic1[a2[k][0]]
    font=run.font
    font.size = Pt(16)
    font.bold = True
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_2

    p = tf.add_paragraph()
    run = p.add_run()
    run.text = min1 + "-" + max1

    shape = slide.shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_1, Inches(0.1), Inches(1.05) ,Inches(7),Inches(0.02))
    line = shape.line
    line.fill.background()

    font = run.font
    font.name = 'Calibri'
    font.size = Pt(16)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_1



    chart_data = ChartData()
    chart_data.categories = b1
    chart_data.add_series(a2[k][0],row)
    chart_data.add_series(a2[k][0],'0')
    x, y, cx, cy = Inches(0), Inches(1.5), Inches(10), Inches(5.75)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = False
    chart.series[0].smooth = True
    # value_series=chart.series[0]
    # value_series.fill.solid()
    # value_series.fill.fore_color.theme_color=MSO_THEME_COLOR.ACCENT_2
    value_axis = chart.value_axis
    tick_labels = value_axis.tick_labels
    tick_labels.number_format = '"$"0.00'
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(12)
    category_axis=chart.category_axis
    tick_labels = category_axis.tick_labels
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(12)
    category_axis = chart.category_axis
    category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW


##################################################        FHFA                         ######################################

# title_slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(title_slide_layout)
# title = slide.shapes.title
# subtitle = slide.placeholders[1]

regionNames={"VA":"Virginia","VB-Norfolk":"Virginia Beach","Washington-Alexandria":"Washington","Blacksburg":"Blacksburg","Charlottesville":"Charlottesville",
			"Winchester":"Winchester","Harrisonburg":"Harrisonburg","Lynchburg":"Lynchburg","Richmond":"Richmond","Roanoke":"Roanoke","Staunton-Waynesboro":"Staunton"};

# title.text = "FHFA"
# subtitle.text = "FHFA Housing price Index by Quarter"
f4 = open("./Csv/FHFA_Quarterly.csv",'r')
csv_f4 = csv.reader(f4)
a5=[]
for row in csv_f4:
    a5.append(row)
a5=zip(*a5)
max1 = 0
min1 = 0
min1 = a5[0][1].split('/')[-1]
max1 = a5[0][-1].split('/')[-1]
b1=[]
for m in range(1,len(a5[0])):
    b1.append(a5[0][m])

for l in range(1,len(a5)):

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(header_img, Inches(0), Inches(0))
    chart_data = ChartData()
    chart_data.categories = b1
    row=[]
    for k in range(1,len(a5[l])):
        row.append(float(a5[l][k]))
    chart_data.add_series(a5[l][0],row)

    txBox = slide.shapes.add_textbox(Inches(0), Inches(0.75) ,Inches(10),Inches(0.5))
    tf = txBox.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    run = p1.add_run()
    run.text ="FHFA Housing price Index by Quarter for " + regionNames[a5[l][0]]
    font=run.font
    font.size = Pt(16)
    font.bold = True
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_2

    p = tf.add_paragraph()
    run = p.add_run()
    run.text = min1 + "-" + max1

    shape = slide.shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_1, Inches(0.1), Inches(1.05) ,Inches(7),Inches(0.02))
    line = shape.line
    line.fill.background()

    font = run.font
    font.name = 'Calibri'
    font.size = Pt(16)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    #chart_data.add_series('Series 1', (19.2, 21.4, 16.7))

    # add chart to slide --------------------
    x, y, cx, cy = Inches(0), Inches(1.5), Inches(10), Inches(6)
    chart=slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    value_series=chart.series[0]
    value_series.fill.solid()
    value_series.fill.fore_color.theme_color=MSO_THEME_COLOR.ACCENT_5
    chart.has_legend = False
    chart.series[0].smooth = True
    value_axis = chart.value_axis
    tick_labels = value_axis.tick_labels
    tick_labels.number_format = '#,###'
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(11)
    category_axis=chart.category_axis
    tick_labels = category_axis.tick_labels
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(11)
    chart.value_axis.minimum_scale = 0




##################################################      Taxable Retail Sales by Month                      ######################################

# title_slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(title_slide_layout)
# title = slide.shapes.title
# subtitle = slide.placeholders[1]

regionNames={"Virginia":"Virginia",
"Virginia Beach-Norfolk-Newport News, VA-NC Metropolitan Statistical Area":"VirginiaBeach",
"Washington-Arlington-Alexandria, DC-VA-MD-WV Metropolitan Statistical Area":"Washington",
"Blacksburg-Christiansburg-Radford, VA Metropolitan Statistical Area":"Blacksburg",
"Charlottesville, VA Metropolitan Statistical Area":"Charlottesville",
"Winchester, VA-WV Metropolitan Statistical Area":"Winchester",
"Harrisonburg, VA Metropolitan Statistical Area":"Harrisonburg",
"Lynchburg, VA Metropolitan Statistical Area":"Lynchburg",
"Richmond, VA Metropolitan Statistical Area":"Richmond",
"Roanoke, VA Metropolitan Statistical Area":"Roanoke",
"Staunton-Waynesboro, VA Metropolitan Statistical Area":"Staunton",
"Harrisonburg, VA Metropolitan Statistical Area":"Harrisonburg"};

# title.text = "Taxable Retail Sales by Month"
# subtitle.text = "(In Millions of Dollars)"


f5 = open("./Csv/local_option_sales_data.csv",'r')
csv_f5 = csv.reader(f5)
a5=[]
for row in csv_f5:
    a5.append(row)
a5=zip(*a5)
max1 = 0
min1 = 0
min1 = a5[0][1].split('/')[-1]
max1 = a5[0][-1].split('/')[-1]


for k in range(1,len(a5)):

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(header_img, Inches(0), Inches(0))
    chart_data = ChartData()
    b1=[]
    for m in range(0,len(a5[0])):
        if m!=0:
            b1.append(a5[0][m])
    chart_data.categories = b1
    row=[]
    for l in range(1,len(a5[k])):
        #print b[k]
        row.append(float(a5[k][l])/10000)
    oo = a5[k][0]
    #title.text = oo

    txBox = slide.shapes.add_textbox(Inches(0), Inches(0.75) ,Inches(10),Inches(0.5))
    tf = txBox.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    run = p1.add_run()
    run.text ="Taxable Retail Sales by Month for " + regionNames[oo] + " (In Millions of Dollars)"
    domainStart = { "Virginia" : 4000, "VirginiaBeach" : 1000,   
        "Blacksburg" : 80,   
        "Charlottesville" : 100,    
        "Harrisonburg" : 80,   
        "Lynchburg" : 100,   
        "Richmond" : 800,  
        "Roanoke" : 200, 
        "Staunton" : 60, 
        "Winchester" : 50, 
        "Washington" : 2000 }
    font=run.font
    font.size = Pt(16)
    font.bold = True
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_2

    p = tf.add_paragraph()
    run = p.add_run()
    run.text = min1 + "-" + max1

    shape = slide.shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_1, Inches(0.1), Inches(1.05) ,Inches(7),Inches(0.02))
    line = shape.line
    line.fill.background()

    font = run.font
    font.name = 'Calibri'
    font.size = Pt(16)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_1




    chart_data.add_series(a5[k][0],row)
    chart_data.add_series(a5[k][0],'0')
    x, y, cx, cy = Inches(0), Inches(1.5), Inches(10), Inches(5.75)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = False
    chart.series[0].smooth = True
    value_axis = chart.value_axis
    tick_labels = value_axis.tick_labels
    tick_labels.number_format = '"$"#,##0.00'
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(12)
    category_axis=chart.category_axis
    tick_labels = category_axis.tick_labels
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(12)
    chart.value_axis.minimum_scale = domainStart[regionNames[oo]]


##################################################           uiclaims                ######################################

# title_slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(title_slide_layout)
# title = slide.shapes.title
# subtitle = slide.placeholders[1]

regionNames={"Virginia":"Virginia",
"Virginia Beach-Norfolk-Newport News, VA-NC Metropolitan Statistical Area, VA part":"Virginia Beach",
"Washington-Arlington-Alexandria, DC-VA-MD-WV Metropolitan Statistical Area, VA part":"Washington",
"Blacksburg-Christiansburg-Radford, VA Metropolitan Statistical Area":"Blacksburg",
"Charlottesville, VA Metropolitan Statistical Area":"Charlottesville",
"Winchester, VA-WV Metropolitan Statistical Area, VA part":"Winchester",
"Harrisonburg, VA Metropolitan Statistical Area":"Harrisonburg",
"Lynchburg, VA Metropolitan Statistical Area":"Lynchburg",
"Richmond, VA Metropolitan Statistical Area":"Richmond",
"Roanoke, VA Metropolitan Statistical Area":"Roanoke",
"Staunton-Waynesboro, VA Metropolitan Statistical Area":"Staunton",
"Harrisonburg, VA Metropolitan Statistical Area":"Harrisonburg"};


# title.text = "Initial Unemployment Claims by Month"
# subtitle.text = " "
f6 = open("./Csv/uiclaims.csv",'r')
csv_f6 = csv.reader(f6)
a5=[]
for row in csv_f6:
    a5.append(row)
a5=zip(*a5)

max1 = 0
min1 = 0
min1 = a5[0][1].split('/')[-1]
max1 = a5[0][-1].split('/')[-1]


for k in range(1,len(a5)):

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(header_img, Inches(0), Inches(0))
    chart_data = ChartData()
    b1=[]
    for m in range(0,len(a5[0])):
        if m!=0:
            b1.append(a5[0][m])
    chart_data.categories = b1
    row=[]
    for l in range(1,len(a5[k])):
        #print b[k]
        row.append(float(a5[k][l]))
    oo = a5[k][0]
    txBox = slide.shapes.add_textbox(Inches(0), Inches(0.75) ,Inches(10),Inches(0.5))
    tf = txBox.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    run = p1.add_run()
    #print oo
    run.text ="Initial Unemployment Claims by Month for " + regionNames[oo]
    font=run.font
    font.size = Pt(16)
    font.bold = True
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_2

    p = tf.add_paragraph()
    run = p.add_run()
    run.text = min1 + "-" + max1

    shape = slide.shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_1, Inches(0.1), Inches(1.05) ,Inches(7),Inches(0.02))
    line = shape.line
    line.fill.background()

    font = run.font
    font.name = 'Calibri'
    font.size = Pt(16)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_1


    chart_data.add_series(a5[k][0],row)
    chart_data.add_series(a5[k][0],'0')
    x, y, cx, cy = Inches(0), Inches(1.5), Inches(10), Inches(5.75)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = False
    chart.series[0].smooth = True
    value_axis = chart.value_axis
    tick_labels = value_axis.tick_labels
    tick_labels.number_format = '#,##0'
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(12)
    category_axis=chart.category_axis
    tick_labels = category_axis.tick_labels
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(12)
    chart.value_axis.minimum_scale = 0


##################################################        Average weekly wages                         ######################################

# title_slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(title_slide_layout)
# title = slide.shapes.title
# subtitle = slide.placeholders[1]

regionNames={"Virginia":"Virginia","Hampton Roads":"Virginia Beach","Wash DC":"Washington","Blacksburg":"Blacksburg","Charlottesville":"Charlottesville",
			"Winchestor":"Winchester","Harrisonburg":"Harrisonburg","Lynchburg":"Lynchburg","Richmond":"Richmond","Roanoke":"Roanoke","Staunton":"Staunton"};

# title.text = "Average Weekly Wages by Quarter"
# subtitle.text = " "
f7 = open("./Csv/WEEKLY_WAGES_Quarterly.csv",'r')
csv_f7 = csv.reader(f7)
a5=[]

for row in csv_f7:
    #print row
    a5.append(row)
a5=zip(*a5)
max1 = 0
min1 = 0
#print a5
min1 = a5[0][1].split('/')[-1]
max1 = a5[0][-1].split('/')[-1]
b1=[]
for m in range(1,len(a5[0])):
    b1.append(a5[0][m])

for l in range(1,len(a5)):

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(header_img, Inches(0), Inches(0))
    row=[]
    # b1=[]
    # for m in range(0,len(a5[0])):
    #     if m!=0:
    #         b1.append(a5[0][m])
    # print b1
    for k in range(1,len(a5[l])):
        row.append(float(a5[l][k]))
        #print row
    #print a5[l][0]

    oo = a5[l][0]
    #chart_data.add_series('Series 1', (19.2, 21.4, 16.7))

    txBox = slide.shapes.add_textbox(Inches(0), Inches(0.75) ,Inches(10),Inches(0.5))
    tf = txBox.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    run = p1.add_run()
    run.text ="Average Weekly Wages by Quarter for " + regionNames[oo]
    font=run.font
    font.size = Pt(16)
    font.bold = True
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_2

    p = tf.add_paragraph()
    run = p.add_run()
    run.text = min1 + "-" + max1

    shape = slide.shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_1, Inches(0.1), Inches(1.05) ,Inches(7),Inches(0.02))
    line = shape.line
    line.fill.background()

    font = run.font
    font.name = 'Calibri'
    font.size = Pt(16)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_1


    # add chart to slide --------------------
    chart_data = ChartData()
    chart_data.categories = b1
    chart_data.add_series(a5[l][0],row)
    x, y, cx, cy = Inches(0), Inches(1.5), Inches(10), Inches(5.75)
    chart=slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    value_series=chart.series[0]
    value_series.fill.solid()
    value_series.fill.fore_color.theme_color=MSO_THEME_COLOR.ACCENT_5
    chart.has_legend = False
    chart.series[0].smooth = True
    value_axis = chart.value_axis
    tick_labels = value_axis.tick_labels
    tick_labels.number_format = '"$"#,##0.00'
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(12)
    category_axis=chart.category_axis
    tick_labels = category_axis.tick_labels
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(12)
    chart.value_axis.minimum_scale = 0

######################################         Seasonally Adjusted Size of Labor Force by Month(In Thousands)         ########################################

# title_slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(title_slide_layout)
# title = slide.shapes.title
# subtitle = slide.placeholders[1]

regionNames={"VALF":"Virginia","VIRG251LF":"VirginiaBeach","WASH911LF":"Washington","BLAC951LF":"Blacksburg","CHAR851LF":"Charlottesville",
		"WINC051LF":"Winchester","HARR551LF":"Harrisonburg","LYNC351LF":"Lynchburg","RICH051LF":"Richmond","ROAN251LF":"Roanoke","Staunton":"","HARR551LF":"Harrisonburg"};
# title.text = "Seasonally Adjusted Size of Labor Force by Month"
# subtitle.text = "(In Thousands)"

f9 = open('./Csv/LABOR_FORCE_lf_ssa_Monthly.csv')
csv_f9 = csv.reader(f9)
a1=[]

for row in csv_f9:
  a1.append(row)
a1=zip(*a1)
#print a1
max1 = 0
min1 = 0
#print a5
min1 = a1[0][1].split('/')[-1]
max1 = a1[0][-1].split('/')[-1]



for k in range(1,len(a1)):
    b1=[]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(header_img, Inches(0), Inches(0))
    chart_data = ChartData()
    b=[]
    for m in range(0,len(a1[0])):
        if m!=0:
            b1.append(a1[0][m])
    chart_data.categories = b1
    row=[]
    for l in range(1,len(a1[k])):
        
        if float(a1[k][l]) == 0 :
            #print float(a1[k][l]), b1[l-1]
            del(b1[l-1])
        else:
            row.append(float(a1[k][l])/1000)
    oo = a1[k][0]
    txBox = slide.shapes.add_textbox(Inches(0), Inches(0.75) ,Inches(10),Inches(0.5))
    tf = txBox.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    run = p1.add_run()
    run.text ="Seasonally Adjusted Size of Labor Force by Month for " + regionNames[oo] + " (in Thousands)"
    domainStart = { "Virginia" : 3400, "VirginiaBeach" : 700,   
        "Blacksburg" : 65,   
        "Charlottesville" : 90,    
        "Harrisonburg" : 54,   
        "Lynchburg" : 90,   
        "Richmond" : 540,  
        "Roanoke" : 120, 
        "Staunton" : 40, 
        "Winchester" : 50, 
        "Washington" : 2400 }
    font=run.font
    font.size = Pt(16)
    font.bold = True
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_2

    p = tf.add_paragraph()
    run = p.add_run()
    run.text = min1 + "-" + max1

    shape = slide.shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_1, Inches(0.1), Inches(1.05) ,Inches(7),Inches(0.02))
    line = shape.line
    line.fill.background()

    font = run.font
    font.name = 'Calibri'
    font.size = Pt(16)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_1



    chart_data.add_series(a1[k][0],row)
    chart_data.add_series(a1[k][0],'0')
    x, y, cx, cy = Inches(0), Inches(1.5), Inches(10), Inches(5.75)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = False
    chart.series[0].smooth = True
    value_axis = chart.value_axis
    tick_labels = value_axis.tick_labels
    tick_labels.number_format = '#,###'
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(12)
    category_axis=chart.category_axis
    category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
    tick_labels = category_axis.tick_labels
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(12)
    chart.value_axis.minimum_scale = domainStart[regionNames[oo]]




######################################         Not Seasonally Adjusted Size of Labor Force by Month(In Thousands)     ########################################

# title_slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(title_slide_layout)
# title = slide.shapes.title
# subtitle = slide.placeholders[1]
regionNames={"VALFN":"Virginia","VIRG251LFN":"VirginiaBeach","WASH911LFN":"Washington","BLAC951LFN":"Blacksburg","CHAR851LFN":"Charlottesville",
			"WINC051LFN":"Winchester","HARR551LFN":"Harrisonburg","LYNC351LFN":"Lynchburg","RICH051LFN":"Richmond","ROAN251LFN":"Roanoke","staunton":"Staunton","HARR551LFN":"Harrisonburg"};
# title.text = "Not Seasonally Adjusted Size of Labor Force by Month"
# subtitle.text = "(In Thousands)"

f8 = open('./Csv/LABOR_FORCE_lf_nsa_Monthly2.csv')
csv_f8 = csv.reader(f8)
a1=[]

for row in csv_f8:
  a1.append(row)
a1=zip(*a1)
max1 = 0
min1 = 0
#print a5
min1 = a1[0][1].split('/')[-1]
max1 = a1[0][-1].split('/')[-1]



for k in range(1,len(a1)):
    b1=[]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(header_img, Inches(0), Inches(0))
    chart_data = ChartData()
    b=[]
    for m in range(0,len(a1[0])):
        if m!=0:
            b1.append(a1[0][m])
    chart_data.categories = b1
    row=[]
    for l in range(1,len(a1[k])):
        # if l==len(a1[k]) or l==len(a1[k])-1 or l==len(a1[k])-2:
        #     if float(a1[k][l])!=0:
        #         row.append(float(a1[k][l])/1000)
        # else:
            if float(a1[k][l]) == 0 :
                #print float(a1[k][l]), b1[l-1]
                del(b1[l-1])
            else:
                row.append(float(a1[k][l])/1000)


    # for l in range(1,len(a1[k])):
    #     row.append(float(a1[k][l])/1000)
    oo = a1[k][0]
    txBox = slide.shapes.add_textbox(Inches(0), Inches(0.75) ,Inches(10),Inches(0.5))
    tf = txBox.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    run = p1.add_run()
    run.text ="Not Seasonally Adjusted Size of Labor Force by Month for " + regionNames[oo]
    domainStart = { "Virginia" : 3400, "VirginiaBeach" : 700,   
        "Blacksburg" : 65,   
        "Charlottesville" : 90,    
        "Harrisonburg" : 54,   
        "Lynchburg" : 90,   
        "Richmond" : 540,  
        "Roanoke" : 120, 
        "Staunton" : 40, 
        "Winchester" : 50, 
        "Washington" : 2400 }
	#domainStart = { "Virginia" : 3400, "VirginiaBeach" : 700, "Blacksburg" : 65, "Charlottesville" : 90, "Harrisonburg" : 54, "Lynchburg" : 90, "Richmond" : 540, "Roanoke" : 120, "Staunton" : 40, "Winchester" : 50, "Washington" : 2400}                        
    font=run.font
    font.size = Pt(16)
    font.bold = True
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_2

    p = tf.add_paragraph()
    run = p.add_run()
    run.text = min1 + "-" + max1

    shape = slide.shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_1, Inches(0.1), Inches(1.05) ,Inches(7),Inches(0.02))
    line = shape.line
    line.fill.background()

    font = run.font
    font.name = 'Calibri'
    font.size = Pt(16)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_1




    chart_data.add_series(a1[k][0],row)
    chart_data.add_series(a1[k][0],'0')
    x, y, cx, cy = Inches(0), Inches(1.5), Inches(10), Inches(5.75)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = False
    chart.series[0].smooth = True
    value_axis = chart.value_axis
    tick_labels = value_axis.tick_labels
    tick_labels.number_format = '#,###.00'
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(12)
    category_axis=chart.category_axis
    tick_labels = category_axis.tick_labels
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(12)
    chart.value_axis.minimum_scale = domainStart[regionNames[oo]]




######################################         Unemployment ssa         ########################################

# title_slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(title_slide_layout)
# title = slide.shapes.title
# subtitle = slide.placeholders[1]
#
# title.text = "Seasonally Adjusted Unemployment Rate by Month"
# subtitle.text = " "

f10 = open('./Csv/labor_force_data.csv')
csv_f10 = csv.reader(f10)
a1=[]

regionNames={"VAUR":"Virginia","VIRG251UR":"Virginia Beach","WASH911UR":"Washington","BLAC951UR":"Blacksburg","CHAR851UR":"Charlottesville",
			"WINC051UR":"Winchester","HARR551UR":"Harrisonburg","LYNC351UR":"Lynchburg","RICH051UR":"Richmond","ROAN251UR":"Roanoke"};

for row in csv_f10:
  a1.append(row)
a1=zip(*a1)
#print a1
max1 = 0
min1 = 0
#print a5
min1 = a1[0][1].split('/')[-1]
max1 = a1[0][-1].split('/')[-1]



b1=[]
for k in range(1,len(a1)):

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(header_img, Inches(0), Inches(0))
    chart_data = ChartData()
    b=[]
    for m in range(0,len(a1[0])):
        if m!=0:
            b1.append(a1[0][m])
    chart_data.categories = b1
    row=[]
    for l in range(1,len(a1[k])):

            if float(a1[k][l]) == 0 :   
                #print float(a1[k][l]), b1[l-1]
                del(b1[l-1])
            else:
                row.append(float(a1[k][l]))

        #row.append(float(a1[k][l]))
    oo = a1[k][0]
    txBox = slide.shapes.add_textbox(Inches(0), Inches(0.75) ,Inches(10),Inches(0.5))
    tf = txBox.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    run = p1.add_run()
    run.text ="Seasonally Adjusted Unemployment Rate by Month for " + regionNames[oo]
    font=run.font
    font.size = Pt(16)
    font.bold = True
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_2

    p = tf.add_paragraph()
    run = p.add_run()
    run.text = min1 + "-" + max1

    shape = slide.shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_1, Inches(0.1), Inches(1.05) ,Inches(7),Inches(0.02))
    line = shape.line
    line.fill.background()

    font = run.font
    font.name = 'Calibri'
    font.size = Pt(16)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_1


    chart_data.add_series(a1[k][0],row)
    chart_data.add_series(a1[k][0],'0')
    x, y, cx, cy = Inches(0), Inches(1.5), Inches(10), Inches(5.75)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = False
    chart.series[0].smooth = True
    value_axis = chart.value_axis
    tick_labels = value_axis.tick_labels
    chart.value_axis.minimum_scale = 2
    tick_labels.number_format = '0.0"%"'
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(12)
    category_axis=chart.category_axis
    tick_labels = category_axis.tick_labels
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(12)


######################################         Unemployment nsa     ########################################

# title_slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(title_slide_layout)
# title = slide.shapes.title
# subtitle = slide.placeholders[1]
#
# title.text = "Not Seasonally Adjusted Unemployment Rate by Month"
# subtitle.text = " "

regionNames={"VAURN":"Virginia","VIRG251URN":"Virginia Beach","WASH911URN":"Washington","BLAC951URN":"Blacksburg","CHAR851URN":"Charlottesville",
			"WINC051URN":"Winchester","HARR551URN":"Harrisonburg","LYNC351URN":"Lynchburg","RICH051URN":"Richmond","ROAN251URN":"Roanoke","staunton":"Staunton"};

f11 = open('./Csv/labor_force_dataNsa.csv')
csv_f11 = csv.reader(f11)
a1=[]

for row in csv_f11:
  a1.append(row)
a1=zip(*a1)
#print a1
max1 = 0
min1 = 0
#print a5
min1 = a1[0][1].split('/')[-1]
max1 = a1[0][-1].split('/')[-1]


for k in range(1,len(a1)):
    b1=[]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(header_img, Inches(0), Inches(0))
    chart_data = ChartData()
    b=[]
    for m in range(0,len(a1[0])):
        if m!=0:
            b1.append(a1[0][m])
    chart_data.categories = b1
    row=[]
    for l in range(1,len(a1[k])):
        if float(a1[k][l]) == 0 :
                #print float(a1[k][l]), b1[l-1]
            del(b1[l-1])
        else:
            row.append(float(a1[k][l]))
        #row.append(float(a1[k][l]))
    oo = a1[k][0]
    txBox = slide.shapes.add_textbox(Inches(0), Inches(0.75) ,Inches(10),Inches(0.5))
    tf = txBox.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    run = p1.add_run()
    run.text ="Not Seasonally Adjusted Unemployment Rate by Month for " + regionNames[oo]
    font=run.font
    font.size = Pt(16)
    font.bold = True
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_2

    p = tf.add_paragraph()
    run = p.add_run()
    run.text = min1 + "-" + max1

    shape = slide.shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_1, Inches(0.1), Inches(1.05) ,Inches(7),Inches(0.02))
    line = shape.line
    line.fill.background()

    font = run.font
    font.name = 'Calibri'
    font.size = Pt(16)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_1


    chart_data.add_series(a1[k][0],row)
    chart_data.add_series(a1[k][0],'0')
    x, y, cx, cy = Inches(0), Inches(1.5), Inches(10), Inches(5.75)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = False
    chart.series[0].smooth = True
    value_axis = chart.value_axis
    tick_labels = value_axis.tick_labels
    chart.value_axis.minimum_scale = 2
    tick_labels.number_format = '0.0"%"'
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(12)
    category_axis=chart.category_axis
    tick_labels = category_axis.tick_labels
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(12)





prs.save('CEAP.pptx')
